import logging
import os
import uuid
import re
from datetime import datetime
from pathlib import Path
from typing import List, Dict, Tuple

import docx
import pytesseract
from PIL import Image
from pypdf import PdfReader
from pptx import Presentation
from pdf2image import convert_from_path

from sentence_transformers import SentenceTransformer

from qdrant_client import QdrantClient
from qdrant_client.http import models as rest_models

from app.utils.logger import log_service
from qdrant_client import QdrantClient
from qdrant_client.http import models as rest_models
from qdrant_client.http.exceptions import UnexpectedResponse

# -------------------------
# CONFIG
# -------------------------
DEFAULT_QDRANT_HOST = os.getenv("QDRANT_HOST", "vector_db")
DEFAULT_QDRANT_REST_PORT = int(os.getenv("QDRANT_REST_PORT", 6333))
DEFAULT_QDRANT_GRPC_PORT = int(os.getenv("QDRANT_GRPC_PORT", 6334))

BATCH_SIZE = 32          # MiniLM is fast; 32 is ideal for CPU
CHUNK_SIZE = 800
OVERLAP = 100

TEMP_DIR = Path("temp_uploads")
TEMP_DIR.mkdir(exist_ok=True)


# ====================================
# TEXT HELPERS
# ====================================
def clean_text(s: str) -> str:
    s = s.replace("\r", " ")
    s = re.sub(r"\n{2,}", "\n\n", s)
    s = re.sub(r"\s{2,}", " ", s)
    return s.strip()

def split_text_to_chunks(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = OVERLAP) -> List[str]:
    if overlap >= chunk_size:
        raise ValueError("overlap must be smaller than chunk_size")
    tokens = text.split()
    chunks = []
    i = 0
    n = len(tokens)
    while i < n:
        j = min(i + chunk_size, n)
        chunks.append(" ".join(tokens[i:j]))
        next_i = j - overlap
        if next_i <= i:
            next_i = i + 1
        i = next_i
    return chunks


# ====================================
# FILE EXTRACTION
# ====================================
def extract_text_from_pdf(path: str):
    pages = []
    reader = PdfReader(path)
    for page_num, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            images = convert_from_path(path, dpi=200, first_page=page_num, last_page=page_num)
            text = pytesseract.image_to_string(images[0])
        pages.append((page_num, text))
    return pages

def extract_text_from_docx(path: str) -> str:
    doc = docx.Document(path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(parts)


def extract_text_from_pptx(path: str) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = shape.text.strip()
                if t:
                    texts.append(t)
    return "\n\n".join(texts)

def extract_text_from_image(path: str) -> str:
    return pytesseract.image_to_string(Image.open(path))


def extract_text(path: str):
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return extract_text_from_pdf(path)
    if ext == ".docx":
        return extract_text_from_docx(path)
    if ext == ".pptx":
        return extract_text_from_pptx(path)
    if ext in [".txt", ".md"]:
        return Path(path).read_text(encoding="utf-8", errors="ignore")
    if ext in [".png", ".jpg", ".jpeg", ".bmp", ".tiff"]:
        return extract_text_from_image(path)
    raise ValueError(f"Unsupported extension: {ext}")


# ====================================
# EMBEDDING — MiniLM Version
# ====================================

class Embedder:

    def __init__(self, model_name="all-MiniLM-L6-v2"):
        self.model_name = model_name

        # Inside Docker: app/app/embedding_models/<model_name>
        base_dir = Path(__file__).resolve().parent.parent
        model_dir = base_dir / "embedding_models" / model_name

        print(f"[Embedder] Loading local model from: {model_dir}")

        self.model = SentenceTransformer(str(model_dir), device="cpu")
        self.dim = self.model.get_sentence_embedding_dimension()
        self.is_bge = False

        print(f"[Embedder] Loaded. Dim = {self.dim}")


    def embed(self, texts: List[str]):
        return self.model.encode(
            texts,
            show_progress_bar=False,
            convert_to_numpy=True
        )


_embedder_instance = None

@log_service
def get_embedder():
    global _embedder_instance
    if _embedder_instance is None:
        _embedder_instance = Embedder()
    return _embedder_instance


# ====================================
# QDRANT RETRIEVAL
# ====================================
from qdrant_client.http import models as rest_models

def retrieve(
    query: str,
    embedder: Embedder,
    qc: QdrantClient,
    collection: str,
    top_k: int = 5,
    filters=None,
):
    q_vec = embedder.model.encode(
        [query],
        show_progress_bar=False,
        convert_to_numpy=True,
    )[0]

    query_request = rest_models.QueryRequest(
        query=q_vec,
        limit=top_k,
        filter=filters,
        with_payload=True,
    )

    batch_results = qc.query_batch_points(
        collection_name=collection,
        requests=[query_request],
    )

    # ✅ Extract actual points
    response = batch_results[0] if batch_results else None
    hits = response.points if response else []

    formatted = []

    for hit in hits:
        # hit is a ScoredPoint-like object OR tuple depending on backend
        # Normalize defensively
        if hasattr(hit, "score"):
            score = hit.score
            payload = hit.payload or {}
        else:
            # tuple variants
            if len(hit) == 3:
                _, score, payload = hit
            elif len(hit) == 2:
                _, score = hit
                payload = {}
            else:
                continue

        if isinstance(score, list):
            score = score[0] if score else 0.0

        formatted.append(
            f"[score={float(score):.4f} | "
            f"page {payload.get('page')} | "
            f"chunk {payload.get('chunk_index')}]\n"
            f"{payload.get('text_preview', '')}"
        )

    return hits, "\n\n---\n\n".join(formatted)




# ====================================
# QDRANT HELPERS
# ====================================
def ensure_collection(
    qc: QdrantClient,
    collection_name: str,
    vector_size: int,
):
    collections = qc.get_collections().collections
    names = {c.name for c in collections}

    if collection_name in names:
        return

    print(f"[Qdrant] Creating collection '{collection_name}' (dim={vector_size})")

    qc.create_collection(
        collection_name=collection_name,
        vectors_config=rest_models.VectorParams(
            size=vector_size,
            distance=rest_models.Distance.COSINE,
        ),
    )


def upsert_batches(
    qc: QdrantClient,
    collection_name: str,
    vectors_with_payloads,
    batch_size=BATCH_SIZE,
):
    for i in range(0, len(vectors_with_payloads), batch_size):
        chunk = vectors_with_payloads[i:i + batch_size]

        points = [
            rest_models.PointStruct(
                id=item[0],
                vector=item[1].tolist(),
                payload=item[2],
            )
            for item in chunk
        ]

        try:
            qc.upsert(
                collection_name=collection_name,
                points=points,
                wait=True,  # <-- IMPORTANT
            )
        except UnexpectedResponse as e:
            if e.status_code == 404:
                raise RuntimeError(
                    "Vector collection not initialized"
                ) from e
            raise


_qdrant_client = None


def get_qdrant_client():
    global _qdrant_client
    if _qdrant_client is None:
        _qdrant_client = QdrantClient(
            host=DEFAULT_QDRANT_HOST,
            port=DEFAULT_QDRANT_REST_PORT,
            grpc_port=DEFAULT_QDRANT_GRPC_PORT,
            prefer_grpc=False
        )
    return _qdrant_client


# ====================================
# INGESTION PIPELINE
# ====================================
@log_service
def process_file_sync(path: str, embedder: Embedder, chunk_size: int = CHUNK_SIZE, overlap: int = OVERLAP) -> List[dict]:

    ext = Path(path).suffix.lower()
    extracted = extract_text(path)

    chunks = []
    page_numbers = []

    if ext == ".pdf":
        for page_num, text in extracted:
            clean = clean_text(text)
            page_chunks = split_text_to_chunks(clean, chunk_size, overlap)
            chunks.extend(page_chunks)
            page_numbers.extend([page_num] * len(page_chunks))
    else:
        clean = clean_text(extracted)
        page_chunks = split_text_to_chunks(clean, chunk_size, overlap)
        chunks = page_chunks
        page_numbers = [None] * len(page_chunks)

    if not chunks:
        return []

    # ----------------------------
    # EMBED THE CHUNKS (CPU BOUND)
    # ----------------------------
    vectors = embedder.embed(chunks)

    results = []
    for i, chunk in enumerate(chunks):
        results.append(
            {
                "text": chunk,
                "page": page_numbers[i],
                "text_preview": chunk[:300],
                "length_tokens": len(chunk.split()),
                "vector": vectors[i].tolist(),   # <--- now valid
            }
        )

    return results


# =======================================
# MAIN SERVICE FUNCTION
# =======================================
def save_temp_file_sync(file_bytes: bytes, filename: str, temp_dir: str = "temp_uploads") -> str:
    Path(temp_dir).mkdir(parents=True, exist_ok=True)
    ext = Path(filename).suffix
    name = f"{uuid.uuid4().hex}{ext}"
    path = Path(temp_dir) / name
    with open(path, "wb") as f:
        f.write(file_bytes)
    return str(path)

